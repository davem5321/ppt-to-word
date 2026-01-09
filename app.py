import streamlit as st
from pptx import Presentation
from docx import Document
from docx.shared import Inches, Pt
from io import BytesIO
import os
import tempfile
import zipfile

st.set_page_config(page_title="PowerPoint to Word Converter", page_icon="📄")

st.title("📄 PowerPoint to Word Converter")
st.write("Upload a PowerPoint file and convert it to a Word document")

# Important notices
st.info("ℹ️ **Requirements:**\n- This tool only supports modern PowerPoint format (**.pptx**)\n- If you have an older .ppt file, please open it in PowerPoint and save it as .pptx first\n- **Sensitivity labels:** If your file has a Microsoft sensitivity label, it must be set to **General** or lower. Files with higher sensitivity levels (Confidential, etc.) are encrypted and cannot be processed.")

# File uploader
uploaded_file = st.file_uploader("Drop your PowerPoint file here", type=['pptx'])

def convert_ppt_to_word(ppt_file):
    """Convert PowerPoint presentation to Word document"""
    # Load the presentation
    prs = Presentation(ppt_file)
    
    # Create a new Word document
    doc = Document()
    
    # Add a title
    doc.add_heading('PowerPoint Conversion', 0)
    
    def clean_text(text):
        """Remove or replace XML-incompatible characters"""
        if not text:
            return ""
        # Remove NULL bytes and control characters except tab, newline, carriage return
        cleaned = ''.join(char for char in text if ord(char) >= 32 or char in '\t\n\r')
        return cleaned
    
    # Iterate through slides
    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Add slide number as heading
        doc.add_heading(f'Slide {slide_idx}', level=1)
        
        # Extract text from shapes
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text.strip():
                    # Check if it's a title or content
                    if shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = clean_text(paragraph.text.strip())
                            if text:
                                # Add paragraph with appropriate formatting
                                p = doc.add_paragraph(text)
                                # Apply basic formatting based on level
                                if paragraph.level == 0:
                                    p.style = 'List Bullet'
                                else:
                                    p.style = 'List Bullet 2'
                
                # Handle tables
                if shape.has_table:
                    table = shape.table
                    doc_table = doc.add_table(rows=len(table.rows), cols=len(table.columns))
                    doc_table.style = 'Light Grid Accent 1'
                    
                    for i, row in enumerate(table.rows):
                        for j, cell in enumerate(row.cells):
                            doc_table.rows[i].cells[j].text = clean_text(cell.text)
            except Exception as shape_error:
                # Log error but continue with other shapes
                doc.add_paragraph(f"[Error processing shape: {str(shape_error)}]")
        
        # Add a page break after each slide (except the last one)
        if slide_idx < len(prs.slides):
            doc.add_page_break()
    
    # Save to BytesIO object
    doc_buffer = BytesIO()
    doc.save(doc_buffer)
    doc_buffer.seek(0)
    
    return doc_buffer

if uploaded_file is not None:
    try:
        with st.spinner('Converting PowerPoint to Word...'):
            # Display file info
            st.info(f"📄 File: {uploaded_file.name}")
            
            # Reset uploaded file to beginning (important!)
            uploaded_file.seek(0)
            
            # Read the bytes
            file_bytes = uploaded_file.getvalue()
            
            # Create a temporary file path
            tmp_dir = tempfile.gettempdir()
            tmp_filename = f"temp_{uploaded_file.name}"
            tmp_path = os.path.join(tmp_dir, tmp_filename)
            
            try:
                # Write the file using standard file operations
                with open(tmp_path, 'wb') as tmp_file:
                    tmp_file.write(file_bytes)
                    tmp_file.flush()  # Ensure data is written to disk
                    os.fsync(tmp_file.fileno())  # Force write to disk on Windows
                
                # Verify the file was written correctly
                if not os.path.exists(tmp_path):
                    st.error(f"❌ Temp file not created at {tmp_path}")
                    st.stop()
                
                file_size = os.path.getsize(tmp_path)
                st.write(f"✓ Temp file created: {tmp_path}")
                st.write(f"✓ File size: {file_size:,} bytes (original: {len(file_bytes):,} bytes)")
                
                # Verify file is readable
                try:
                    with open(tmp_path, 'rb') as test_file:
                        test_bytes = test_file.read(100)  # Read first 100 bytes
                        st.write(f"✓ File is readable, first bytes: {test_bytes[:20].hex()}")
                except Exception as read_err:
                    st.error(f"❌ Cannot read temp file: {read_err}")
                    st.stop()
                
                # Try to convert the file (let python-pptx do the validation)
                word_file = convert_ppt_to_word(tmp_path)
                
            finally:
                # Clean up temp file
                try:
                    if os.path.exists(tmp_path):
                        os.unlink(tmp_path)
                except:
                    pass
            
            # Success message
            st.success('✅ Conversion completed successfully!')
            
            # Generate filename
            original_filename = uploaded_file.name
            base_filename = os.path.splitext(original_filename)[0]
            word_filename = f"{base_filename}_converted.docx"
            
            # Download button
            st.download_button(
                label="📥 Download Word Document",
                data=word_file,
                file_name=word_filename,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
            
            st.info(f"📝 Your converted document is ready: {word_filename}")
            
    except Exception as e:
        st.error(f"❌ Error during conversion: {str(e)}")
        st.write("Please make sure you uploaded a valid PowerPoint file (.pptx)")
else:
    st.info("👆 Please upload a PowerPoint file to begin")

# Footer
st.markdown("---")
st.markdown("*Built with Streamlit, python-pptx, and python-docx*")
